"""Build .pptx from SlideSpec using python-pptx."""

from __future__ import annotations

import base64
import io
import re
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .models import (
    SLIDE_WIDTH_IN,
    ExportRequest,
    ImageElement,
    ShapeElement,
    SlideSpec,
    TextElement,
    px_rect_to_in,
    slide_height_in,
)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _align(align: str):
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)


def _decode_image_src(src: str) -> bytes:
    if src.startswith("data:"):
        _, b64 = src.split(",", 1)
        return base64.b64decode(b64)
    return base64.b64decode(src)


def _add_text(slide, el: TextElement, slide_w: int, slide_h: int):
    x, y, w, h = px_rect_to_in(el.x, el.y, el.w, el.h, slide_w, slide_h)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = el.text
    p.alignment = _align(el.align)
    run = p.runs[0]
    run.font.size = Pt(el.fontSize or max(8, el.h * 0.75 * 72 / slide_h))
    run.font.bold = el.bold
    run.font.name = "Microsoft YaHei"
    run.font.color.rgb = _hex_to_rgb(el.color)


def _add_shape(slide, el: ShapeElement, slide_w: int, slide_h: int):
    x, y, w, h = px_rect_to_in(el.x, el.y, el.w, el.h, slide_w, slide_h)
    if el.shape == "line":
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(max(h, 0.02))
        )
        line.fill.background()
        if el.stroke:
            line.line.color.rgb = _hex_to_rgb(el.stroke)
            line.line.width = Pt(el.strokeWidth or 1)
        return

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if el.fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(el.fill)
    else:
        shape.fill.background()
    shape.line.fill.background()


def _add_image(slide, el: ImageElement, slide_w: int, slide_h: int):
    x, y, w, h = px_rect_to_in(el.x, el.y, el.w, el.h, slide_w, slide_h)
    data = _decode_image_src(el.src)
    slide.shapes.add_picture(io.BytesIO(data), Inches(x), Inches(y), Inches(w), Inches(h))


def build_presentation(slides: List[SlideSpec]) -> bytes:
    prs = Presentation()
    prs.core_properties.author = "ccbaby"

    for spec in slides:
        sw_in = SLIDE_WIDTH_IN
        sh_in = slide_height_in(spec.width, spec.height)
        prs.slide_width = Inches(sw_in)
        prs.slide_height = Inches(sh_in)

        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        if spec.backgroundImage:
            try:
                data = _decode_image_src(spec.backgroundImage)
                slide.shapes.add_picture(
                    io.BytesIO(data), Inches(0), Inches(0), Inches(sw_in), Inches(sh_in)
                )
            except Exception:
                pass

        for el in spec.elements:
            if isinstance(el, ShapeElement):
                _add_shape(slide, el, spec.width, spec.height)
            elif isinstance(el, ImageElement):
                _add_image(slide, el, spec.width, spec.height)
            elif isinstance(el, TextElement):
                _add_text(slide, el, spec.width, spec.height)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def export_filename(name: str) -> str:
    safe = re.sub(r'[\\/:*?"<>|]', "_", name.strip() or "presentation")
    return safe if safe.endswith(".pptx") else f"{safe}.pptx"
